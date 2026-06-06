"""
Build AI-Powered Medical Equipment Reliability Intelligence Assistant — Architecture PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── Color palette ─────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1e, 0x3a, 0x5f)
BLUE   = RGBColor(0x25, 0x63, 0xeb)
CYAN   = RGBColor(0x06, 0xb6, 0xd4)
GREEN  = RGBColor(0x10, 0xb9, 0x81)
AMBER  = RGBColor(0xf5, 0x9e, 0x0b)
RED    = RGBColor(0xef, 0x44, 0x44)
PURPLE = RGBColor(0x7c, 0x3a, 0xed)
WHITE  = RGBColor(0xff, 0xff, 0xff)
GRAY   = RGBColor(0x64, 0x74, 0x8b)
LGRAY  = RGBColor(0xf1, 0xf5, 0xf9)
DGRAY  = RGBColor(0x1e, 0x29, 0x3b)


# ─── Helpers ───────────────────────────────────────────────────────────────
def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=11, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

def connector(slide, x1, y1, x2, y2, color=GRAY, width=1.5):
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
bg(sl, NAVY)
box(sl, 0, 0, 13.33, 0.08, CYAN)

box(sl, 0.5, 0.9, 0.9, 0.9, CYAN)
txt(sl, "⚕", 0.52, 0.92, 0.86, 0.86, size=30, bold=True, align=PP_ALIGN.CENTER)

txt(sl, "AI-Powered Medical Equipment", 1.65, 0.88, 11, 0.6,
    size=32, bold=True, color=WHITE)
txt(sl, "Reliability Intelligence Assistant", 1.65, 1.46, 11, 0.6,
    size=32, bold=True, color=CYAN)
txt(sl, "Architecture, Data Flow & System Design", 1.65, 2.12, 10, 0.42,
    size=15, color=RGBColor(0xb0, 0xca, 0xe8))

box(sl, 0.5, 2.72, 12.3, 0.04, CYAN)

badges = [
    ("LangGraph Multi-Agent",        BLUE),
    ("GPT-4o mini Guardrails",       PURPLE),
    ("Hybrid RAG (FAISS + BM25)",    GREEN),
    ("DeepEval + LLM-as-Judge",      AMBER),
    ("UCI AI4I 2020 Dataset",        CYAN),
]
bx = 0.5
for label, col in badges:
    box(sl, bx, 2.92, 2.28, 0.38, col)
    txt(sl, label, bx + 0.1, 2.96, 2.08, 0.3, size=9.5, bold=True, align=PP_ALIGN.CENTER)
    bx += 2.46

txt(sl, "Capstone Project  •  Biomedical Engineering AI  •  React + FastAPI + LangGraph",
    0.5, 6.85, 12.3, 0.38, size=9.5, color=RGBColor(0x7a, 0x9a, 0xbf),
    align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — System Architecture Overview
# ══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
bg(sl, LGRAY)
box(sl, 0, 0, 13.33, 0.65, NAVY)
txt(sl, "System Architecture Overview", 0.35, 0.08, 9, 0.48, size=22, bold=True)
txt(sl, "Three-tier architecture: React Frontend → FastAPI Backend → Data Layer",
    0.35, 0.38, 10, 0.28, size=10, color=CYAN)

# --- FRONTEND column ---
box(sl, 0.2, 0.82, 3.6, 6.0, WHITE, BLUE)
box(sl, 0.2, 0.82, 3.6, 0.46, BLUE)
txt(sl, "FRONTEND  (React + Vite)", 0.35, 0.88, 3.3, 0.34, size=10, bold=True)

views = [
    ("📊 Dashboard",          "Stats, charts, failure rates, monthly trends"),
    ("💬 AI Assistant",       "LangGraph chat + guardrail badges + logs"),
    ("📋 Maintenance Records","Paginated, filterable incident table"),
    ("🚨 Anomaly Detection",  "Temp / wear / RPM anomaly analysis"),
    ("🔗 LangGraph Pipeline", "Topology diagram + guardrail tester"),
    ("✅ Evaluation & Judge", "DeepEval metrics + LLM-as-judge scores"),
]
vy = 1.38
for title, desc in views:
    box(sl, 0.35, vy, 3.3, 0.72, LGRAY, RGBColor(0xd1, 0xd5, 0xdb))
    txt(sl, title, 0.48, vy + 0.05, 3.0, 0.28, size=9, bold=True, color=NAVY)
    txt(sl, desc,  0.48, vy + 0.34, 3.0, 0.3,  size=7.5, color=GRAY)
    vy += 0.82

txt(sl, "Axios  /  CORS", 3.83, 3.6, 0.85, 0.28, size=7.5, color=GRAY, align=PP_ALIGN.CENTER)
txt(sl, "→", 3.83, 3.82, 0.85, 0.28, size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# --- BACKEND column ---
box(sl, 4.05, 0.82, 5.2, 6.0, WHITE, NAVY)
box(sl, 4.05, 0.82, 5.2, 0.46, NAVY)
txt(sl, "BACKEND  (FastAPI + Python)", 4.2, 0.88, 4.9, 0.34, size=10, bold=True)

endpoints = [
    ("POST /api/query",               "Full LangGraph pipeline analysis",     BLUE),
    ("GET  /api/equipment/stats",      "Dataset stats + chart data",           GREEN),
    ("GET  /api/maintenance/incidents","Paginated incident records",           GREEN),
    ("POST /api/analyze/anomaly",      "Anomaly detection engine",             AMBER),
    ("POST /api/evaluate",             "DeepEval + LLM-as-judge",             PURPLE),
    ("POST /api/guardrail/validate",   "Standalone guardrail test",            RED),
    ("GET  /api/pipeline/schema",      "LangGraph topology JSON",              CYAN),
]
ey = 1.38
for ep, desc, col in endpoints:
    box(sl, 4.2, ey, 4.95, 0.72, LGRAY, col)
    txt(sl, ep,   4.32, ey + 0.05, 4.7, 0.26, size=8, bold=True, color=col)
    txt(sl, desc, 4.32, ey + 0.32, 4.7, 0.3,  size=7.5, color=GRAY)
    ey += 0.82

txt(sl, "→", 9.28, 3.82, 0.5, 0.28, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# --- DATA column ---
box(sl, 9.55, 0.82, 3.55, 6.0, WHITE, GREEN)
box(sl, 9.55, 0.82, 3.55, 0.46, GREEN)
txt(sl, "DATA LAYER", 9.7, 0.88, 3.25, 0.34, size=10, bold=True)

dstores = [
    ("🗄  ai4i2020.csv",     "UCI AI4I 2020 schema\n10,000 rows × 14 cols",          GREEN),
    ("🔍  FAISS Index",      "IndexFlatIP · 10k × 384-dim\n15 MB on disk",            BLUE),
    ("📝  BM25 Index",       "BM25Okapi keyword index\nAll 10k docs tokenised",        CYAN),
    ("📦  Embeddings",       "all-MiniLM-L6-v2\n10k × 384 float32 .npy",             PURPLE),
    ("⚙   LangGraph State", "MedEquipState TypedDict\n17 fields per pipeline run",    AMBER),
    ("💊  Enriched CSV",     "medical_equipment_maintenance.csv\n32 columns total",    RED),
]
dy = 1.38
for icon, desc, col in dstores:
    box(sl, 9.7, dy, 3.25, 0.82, LGRAY, col)
    txt(sl, icon, 9.82, dy + 0.06, 3.0, 0.28, size=9, bold=True, color=col)
    txt(sl, desc, 9.82, dy + 0.36, 3.0, 0.38, size=7.5, color=GRAY)
    dy += 0.92


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — LangGraph Pipeline Data Flow
# ══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
bg(sl, LGRAY)
box(sl, 0, 0, 13.33, 0.65, NAVY)
txt(sl, "LangGraph Multi-Agent Pipeline — State Machine Data Flow", 0.35, 0.08, 11, 0.48, size=21, bold=True)
txt(sl, "StateGraph · 8 nodes · Conditional routing (VALID/INVALID) · TypedDict state passed between all nodes",
    0.35, 0.38, 12, 0.28, size=9.5, color=CYAN)

# START box
box(sl, 0.15, 2.55, 0.55, 0.5, DGRAY)
txt(sl, "START", 0.15, 2.62, 0.55, 0.36, size=7, bold=True, align=PP_ALIGN.CENTER)

# Main flow nodes
nodes = [
    ("Input\nGuardrail",         1.0,  2.38, 1.55, 0.85, AMBER),
    ("RAG\nRetrieval",           2.9,  2.38, 1.55, 0.85, BLUE),
    ("Retrieval\nAgent",         4.8,  2.38, 1.55, 0.85, GREEN),
    ("Reliability\nAgent",       6.7,  2.38, 1.6,  0.85, GREEN),
    ("Maintenance\nAgent",       8.65, 2.38, 1.6,  0.85, GREEN),
    ("Recommendation\nAgent",   10.6, 2.38, 1.6,  0.85, GREEN),
    ("Output\nGuardrail",       12.55,2.38, 0.65, 0.85, AMBER),
]
for label, x, y, w, h, fill in nodes:
    box(sl, x, y, w, h, fill)
    txt(sl, label, x + 0.08, y + 0.12, w - 0.16, h - 0.24,
        size=8, bold=True, align=PP_ALIGN.CENTER)

# Arrows between nodes
arrow_xs = [0.7, 2.55, 4.45, 6.35, 8.3, 10.25, 12.25]
arrow_y  = 2.8
for ax in arrow_xs:
    connector(sl, ax, arrow_y, ax + 0.35, arrow_y, WHITE, 2)

# END box
box(sl, 13.05, 2.55, 0.23, 0.5, DGRAY)
txt(sl, "END", 13.05, 2.62, 0.23, 0.36, size=6, bold=True, align=PP_ALIGN.CENTER)

# Rejection branch
box(sl, 1.0, 4.15, 1.55, 0.65, RED)
txt(sl, "Rejection\nNode", 1.08, 4.2, 1.4, 0.56, size=8, bold=True, align=PP_ALIGN.CENTER)
box(sl, 1.0, 5.05, 0.8, 0.45, DGRAY)
txt(sl, "END", 1.0, 5.1, 0.8, 0.36, size=7, bold=True, align=PP_ALIGN.CENTER)

connector(sl, 1.77, 3.23, 1.77, 4.15, RED, 1.5)
connector(sl, 1.77, 4.8,  1.4,  5.05, RED, 1.5)
txt(sl, "INVALID", 0.82, 3.55, 0.92, 0.28, size=7.5, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(sl, "VALID",   2.55, 2.5,  0.7,  0.24, size=7.5, bold=True, color=GREEN)

# Detail cards above nodes
details = [
    (1.0,  0.78, 1.55, "GPT-4o mini\nor rule-based\nValidates relevance,\nsafety, intent,\nurgency level",   AMBER),
    (2.9,  0.78, 1.55, "FAISS + BM25\nα=0.6/0.4 fusion\nEmbedding rerank\nMetadata filter\nk×6 candidates", BLUE),
    (4.8,  0.78, 1.55, "Claude Sonnet\nor rule-based\nEquipment focus\nQuery intent\nRelevant incidents",    GREEN),
    (6.7,  0.78, 1.6,  "Claude Sonnet\nor rule-based\nFailure rate %\nReliability 0-100\nAnomaly detection", GREEN),
    (8.65, 0.78, 1.6,  "Claude Sonnet\nor rule-based\nAction plan\nUrgency level\nPreventive steps",         GREEN),
    (10.6, 0.78, 1.6,  "Claude Sonnet\nor rule-based\nFinal synthesis\nRisk + confidence\nFollow-up steps",  GREEN),
    (12.55,0.78, 0.65, "GPT-4o mini\nor rule-based\nQuality 0-100\nSafety rating\nclinical check",           AMBER),
]
for x, y, w, detail, col in details:
    box(sl, x, y, w, 1.45, WHITE, col)
    txt(sl, detail, x + 0.08, y + 0.08, w - 0.16, 1.3, size=7, color=DGRAY)
    connector(sl, x + w/2, y + 1.45, x + w/2, 2.38, RGBColor(0xc0, 0xca, 0xd8), 1)

# State schema strip
box(sl, 0.15, 5.65, 13.05, 1.7, WHITE, NAVY)
txt(sl, "MedEquipState TypedDict — 17 fields", 0.35, 5.72, 8, 0.32, size=9, bold=True, color=NAVY)
fields = [
    "query", "refined_query", "filters", "k", "alpha",
    "input_guardrail_result", "is_valid (router)",
    "retrieved_docs[ ]", "retrieval_result{ }",
    "reliability_result{ }", "maintenance_result{ }",
    "recommendation_result{ }", "output_guardrail_result",
    "error", "analysis_mode", "execution_log[ ]", "node_timings{ }",
]
fx, fy = 0.3, 6.1
for i, f in enumerate(fields):
    col_idx = i % 6
    row_idx = i // 6
    fx2 = 0.3 + col_idx * 2.16
    fy2 = 6.1 + row_idx * 0.46
    box(sl, fx2, fy2, 2.08, 0.38, LGRAY, BLUE)
    txt(sl, f, fx2 + 0.07, fy2 + 0.06, 1.94, 0.28, size=7, color=NAVY)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RAG Architecture
# ══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
bg(sl, LGRAY)
box(sl, 0, 0, 13.33, 0.65, NAVY)
txt(sl, "Hybrid RAG — Retrieval, Fusion & Embedding Reranking", 0.35, 0.08, 11, 0.48, size=21, bold=True)
txt(sl, "Phase 1: FAISS vector search + BM25 keyword search  →  Phase 2: Score fusion (α=0.6/0.4)  →  Phase 3: Embedding rerank  →  Phase 4: Metadata filter",
    0.35, 0.38, 12.6, 0.28, size=9, color=CYAN)

# Phase headers
phases = [
    ("PHASE 1\nVector + Keyword", 0.2,  0.78, 3.35, BLUE),
    ("PHASE 2\nScore Fusion",     3.75, 0.78, 2.85, PURPLE),
    ("PHASE 3\nEmbedding Rerank", 6.8,  0.78, 3.35, RGBColor(0x0d, 0x94, 0x88)),
    ("PHASE 4\nFilter + Budget",  10.35,0.78, 2.75, GREEN),
]
for label, x, y, w, col in phases:
    box(sl, x, y, w, 6.5, WHITE, col)
    box(sl, x, y, w, 0.58, col)
    txt(sl, label, x + 0.12, y + 0.06, w - 0.24, 0.48, size=10, bold=True, align=PP_ALIGN.CENTER)

# Phase 1 content
p1 = [
    ("FAISS Vector Search", BLUE,
     "Model: all-MiniLM-L6-v2\nDim: 384-dimensional\nIndex: IndexFlatIP\n(inner product = cosine\nafter L2 normalization)\nPool: top k×6 candidates\nScore: cosine similarity"),
    ("BM25 Keyword Search", CYAN,
     "Library: rank-bm25 0.2.2\nModel: BM25Okapi\nTokenize: query.lower().split()\nget_scores() across 10k docs\nPool: top k×6 candidates\nScore: BM25 relevance"),
]
py = 1.45
for title, col, detail in p1:
    box(sl, 0.35, py, 3.0, 2.55, LGRAY, col)
    txt(sl, title, 0.48, py + 0.08, 2.75, 0.32, size=9, bold=True, color=col)
    txt(sl, detail, 0.48, py + 0.44, 2.75, 2.05, size=8, color=DGRAY)
    py += 2.72

# Phase 2 content
box(sl, 3.9, 1.45, 2.55, 5.6, LGRAY, PURPLE)
txt(sl, "Score Normalization", 4.02, 1.52, 2.32, 0.3, size=9, bold=True, color=PURPLE)
txt(sl, "min-max normalize both\nvector and BM25 scores\nto [0, 1] range", 4.02, 1.86, 2.32, 0.6, size=8, color=DGRAY)
txt(sl, "Fusion Formula", 4.02, 2.55, 2.32, 0.3, size=9, bold=True, color=PURPLE)
box(sl, 4.02, 2.9, 2.32, 0.9, NAVY)
txt(sl, "combined =\n  0.6 × v_score\n+ 0.4 × b_score",
    4.12, 2.94, 2.12, 0.82, size=8.5, bold=True, color=CYAN)
txt(sl, "Union of candidate sets\nfrom both retrievers\nSorted descending\nby combined score\nPool size: k × 6 docs",
    4.02, 3.92, 2.32, 1.2, size=8, color=DGRAY)

# Phase 3 content
rerank_steps = [
    "① Re-encode query\n   SentenceTransformer\n   → unit vector",
    "② For each candidate\n   cosine_sim =\n   query_emb · doc_emb",
    "③ Metadata bonus\n   +0.08 equip type match\n   +0.06 failure type\n   +0.04 hospital unit\n   +0.05 failure context",
    "④ Recency bonus\n   tool_wear / 10000\n   (max 0.03)",
    "⑤ Final score\n   0.65 × cosine_sim\n   + 0.25 × hybrid\n   + meta + recency",
    "⑥ Sort → top k×4",
]
ry = 1.45
for step in rerank_steps:
    box(sl, 6.95, ry, 3.05, 0.85, LGRAY, RGBColor(0x0d, 0x94, 0x88))
    txt(sl, step, 7.07, ry + 0.07, 2.82, 0.74, size=7.5, color=DGRAY)
    ry += 0.92

# Phase 4 content
filter_items = [
    ("equipment_type",  "Exact match\n(case-insensitive)"),
    ("hospital_unit",   "Exact match"),
    ("severity",        "Low/Medium/High/Critical"),
    ("failure_only",    "machine_failure == 1"),
    ("Token Budget",    "max_tokens = 1200\nword_count × 0.75\nskip over-budget docs"),
    ("Output",          "Top k docs with\nrelevance_score\ndocument_text\n+ all metadata"),
]
fy = 1.45
for fname, fdesc in filter_items:
    box(sl, 10.5, fy, 2.45, 0.88, LGRAY, GREEN)
    txt(sl, fname, 10.62, fy + 0.06, 2.22, 0.26, size=8.5, bold=True, color=GREEN)
    txt(sl, fdesc, 10.62, fy + 0.34, 2.22, 0.5,  size=7.5, color=DGRAY)
    fy += 0.96

# Phase arrows
for ax in [3.55, 6.6, 10.15]:
    txt(sl, "→", ax + 0.02, 3.5, 0.26, 0.36, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Bottom stats bar
box(sl, 0.2, 7.1, 12.9, 0.3, NAVY)
txt(sl, "Dataset: 10,000 records  ·  Embedding: all-MiniLM-L6-v2 (384-dim)  ·  Cache: faiss_index.bin + bm25.pkl + embeddings.npy + documents.pkl  ·  Candidate pool k×6 → rerank → filter → top k",
    0.4, 7.14, 12.5, 0.24, size=7.5, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Multi-Agent System
# ══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
bg(sl, LGRAY)
box(sl, 0, 0, 13.33, 0.65, NAVY)
txt(sl, "Multi-Agent Equipment Intelligence System", 0.35, 0.08, 11, 0.48, size=22, bold=True)
txt(sl, "4 AI Agents · Claude Sonnet 4.6 (with intelligent rule-based fallback) · Sequential execution via LangGraph",
    0.35, 0.38, 12, 0.28, size=9.5, color=CYAN)

agent_data = [
    ("🔍", "Equipment\nRetrieval Agent", BLUE,
     [("Input",   "Query + retrieved docs (up to 8)"),
      ("Output",  "relevant_incidents[]\nquery_intent\nequipment_focus[]\nhospital_units_affected\ntime_period"),
      ("Purpose", "Identifies which equipment types\nand hospital units are most\nrelevant. Filters noise from\nretrieved incident pool."),
      ("Fallback","Keyword matching on equipment\nnames; pattern extraction for\nhospital units and failure types")]),
    ("📈", "Reliability\nAnalysis Agent", GREEN,
     [("Input",   "Query + docs + retrieval_result"),
      ("Output",  "failure_rate (%)\nreliability_score (0–100)\nanomalies[]\nfailure_patterns[]\nseverity_breakdown{}"),
      ("Purpose", "Computes statistical reliability\nmetrics: failure rate, mean tool\nwear, temperature anomaly\ndetection, severity distribution."),
      ("Fallback","Threshold-based anomaly detection;\nstatistical aggregation from\nraw incident data fields")]),
    ("🔧", "Maintenance\nPlanning Agent", AMBER,
     [("Input",   "Query + docs + reliability_result"),
      ("Output",  "immediate_actions[]\nscheduled_maintenance[]\npreventive_measures[]\nmaintenance_urgency"),
      ("Purpose", "Generates prioritised maintenance\nplan with immediate vs scheduled\nvs preventive actions based on\nseverity and failure type."),
      ("Fallback","Rule matrix: failure_type maps\nto standard maintenance protocol;\nurge level from severity field")]),
    ("💡", "Recommendation\nAgent", PURPLE,
     [("Input",   "Query + docs + all 3 prior outputs"),
      ("Output",  "summary\nroot_cause\nconfidence_score (0–100)\nkey_findings[]\naction_plan[]\nrisk_assessment"),
      ("Purpose", "Synthesises all upstream analyses\ninto a final recommendation with\nrisk assessment, root cause,\nand actionable confidence score."),
      ("Fallback","Aggregates rule-based findings;\nscores confidence from doc count\n+ failure count heuristic")]),
]

ax = 0.2
for icon, title, col, items in agent_data:
    box(sl, ax, 0.8, 3.18, 6.5, WHITE, col)
    box(sl, ax, 0.8, 3.18, 0.72, col)
    txt(sl, icon, ax + 0.1, 0.86, 0.52, 0.6, size=22, align=PP_ALIGN.CENTER)
    txt(sl, title, ax + 0.65, 0.88, 2.42, 0.6, size=9.5, bold=True)

    iy = 1.62
    for lbl, val in items:
        txt(sl, lbl + ":", ax + 0.14, iy, 0.88, 0.26, size=7.5, bold=True, color=col)
        txt(sl, val, ax + 0.14, iy + 0.24, 2.9, 0.72, size=7.5, color=DGRAY)
        iy += 1.08

    ax += 3.28

# Orchestration strip
box(sl, 0.2, 7.1, 12.9, 0.3, NAVY)
txt(sl, "LangGraph StateGraph orchestration  ·  Sequential: Retrieval Agent → Reliability Agent → Maintenance Agent → Recommendation Agent  ·  State shared via MedEquipState TypedDict",
    0.35, 7.14, 12.5, 0.24, size=7.5, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Guardrails + DeepEval + LLM-as-Judge
# ══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
bg(sl, LGRAY)
box(sl, 0, 0, 13.33, 0.65, NAVY)
txt(sl, "Quality Assurance — Guardrails + DeepEval + LLM-as-Judge", 0.35, 0.08, 11, 0.48, size=21, bold=True)
txt(sl, "Three-layer QA: Input guardrail → Output guardrail → Post-hoc DeepEval evaluation + LLM-as-judge scoring",
    0.35, 0.38, 12, 0.28, size=9.5, color=CYAN)

# Input Guardrail
box(sl, 0.2, 0.8, 4.0, 2.95, WHITE, AMBER)
box(sl, 0.2, 0.8, 4.0, 0.46, AMBER)
txt(sl, "🛡  INPUT GUARDRAIL", 0.35, 0.86, 3.7, 0.36, size=11, bold=True)
ig = [
    ("Model",    "GPT-4o mini  (rule-based fallback)"),
    ("Checks",   "Medical equipment relevance\nHarmful pattern detection (regex)\nIntent classification\nUrgency: immediate / routine / low"),
    ("Output",   "is_valid → routes pipeline\nrefined_query (improved version)\nintent, equipment_type, urgency"),
    ("Patterns", "hack, exploit, bypass, disable,\nignore safety, override limit"),
]
gy = 1.36
for lbl, val in ig:
    txt(sl, lbl + ":", 0.35, gy, 0.95, 0.28, size=7.5, bold=True, color=AMBER)
    txt(sl, val, 1.32, gy, 2.72, 0.58, size=7.5, color=DGRAY)
    gy += 0.66

# Output Guardrail
box(sl, 4.45, 0.8, 4.0, 2.95, WHITE, AMBER)
box(sl, 4.45, 0.8, 4.0, 0.46, AMBER)
txt(sl, "🛡  OUTPUT GUARDRAIL", 4.6, 0.86, 3.7, 0.36, size=11, bold=True)
og = [
    ("Model",   "GPT-4o mini  (rule-based fallback)"),
    ("Checks",  "Recommendation completeness\nClinical safety validation\nQuality scoring 0–100"),
    ("Scoring", "+10 ≥ 3 action items\n+10 ≥ 3 key findings\n+10 summary ≥ 50 chars\n+10 confidence ≥ 70%\n−20/−30 unsafe language"),
    ("Output",  "quality_score (0–100)\nsafety_rating: safe/caution/unsafe\nis_safe, completeness_check{}"),
]
gy = 1.36
for lbl, val in og:
    txt(sl, lbl + ":", 4.6, gy, 0.88, 0.28, size=7.5, bold=True, color=AMBER)
    txt(sl, val, 5.5, gy, 2.78, 0.75, size=7.5, color=DGRAY)
    gy += 0.66

# DeepEval
box(sl, 8.7, 0.8, 4.4, 2.95, WHITE, PURPLE)
box(sl, 8.7, 0.8, 4.4, 0.46, PURPLE)
txt(sl, "📊  DEEPEVAL RAG METRICS", 8.85, 0.86, 4.1, 0.36, size=11, bold=True)
dv = [
    ("Answer Relevancy",     "Recommendation relevant to query?\nKeyword overlap + entity check"),
    ("Faithfulness",         "Answer grounded in retrieved context?\nEntity traceability scoring"),
    ("Contextual Precision", "Top docs more relevant than lower?\nPrecision@k ordering check"),
    ("Contextual Recall",    "Context covers all answer facts?\nFact extraction + lookup"),
]
dy = 1.36
for metric, desc in dv:
    box(sl, 8.85, dy, 4.1, 0.58, LGRAY, PURPLE)
    txt(sl, metric, 8.97, dy + 0.05, 3.86, 0.24, size=8.5, bold=True, color=PURPLE)
    txt(sl, desc,   8.97, dy + 0.3,  3.86, 0.26, size=7.5, color=DGRAY)
    dy += 0.66

# LLM-as-Judge section
box(sl, 0.2, 3.9, 12.9, 3.38, WHITE, PURPLE)
box(sl, 0.2, 3.9, 12.9, 0.46, PURPLE)
txt(sl, "⚖  LLM-as-Judge — 4-Dimension Rubric  (Total: 0 – 100)", 0.38, 3.96, 12.5, 0.36, size=12, bold=True)

dims = [
    ("Clinical Safety\n(0 – 25)", RED,
     "Are recommendations safe for patients and hospital staff?\n"
     "25: All actions clinically safe; hazardous equipment flagged for shutdown\n"
     "0 : Recommendation could endanger patients (e.g. continue using failed equipment)"),
    ("Technical Accuracy\n(0 – 25)", BLUE,
     "Are technical details correct and evidence-based?\n"
     "25: Temperatures, torque, wear values correctly cited from retrieved incidents\n"
     "0 : Technical claims fabricated or contradict the retrieved incident data"),
    ("Actionability\n(0 – 25)", GREEN,
     "Are recommended actions concrete and immediately implementable?\n"
     "25: Specific, prioritised steps that biomedical engineers can act on now\n"
     "0 : Vague or contradictory actions that cannot be carried out in practice"),
    ("Evidence Basis\n(0 – 25)", AMBER,
     "Is the recommendation grounded in retrieved incidents?\n"
     "25: All claims traceable to specific incidents; no hallucinated data\n"
     "0 : Ignores retrieved incidents entirely or contradicts the evidence"),
]
jx = 0.35
for dim, col, desc in dims:
    box(sl, jx, 4.46, 3.12, 2.68, LGRAY, col)
    txt(sl, dim, jx + 0.12, 4.54, 2.88, 0.52, size=9.5, bold=True, color=col)
    txt(sl, desc, jx + 0.12, 5.08, 2.88, 2.0,  size=7.5, color=DGRAY)
    jx += 3.24

txt(sl, "Judge priority order:  GPT-4o mini  →  Claude Sonnet 4.6  →  Rule-based scoring",
    0.38, 7.1, 12.5, 0.26, size=8, color=GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Dataset & Ingestion Pipeline
# ══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
bg(sl, LGRAY)
box(sl, 0, 0, 13.33, 0.65, NAVY)
txt(sl, "Dataset Architecture — UCI AI4I 2020 → Medical Equipment", 0.35, 0.08, 11, 0.48, size=21, bold=True)
txt(sl, "data_ingestion.py pipeline: Generate ai4i2020.csv (UCI schema) → Enrich with medical context → Build RAG index",
    0.35, 0.38, 12, 0.28, size=9.5, color=CYAN)

# Pipeline flow
flow = [
    ("ai4i2020.csv\n(UCI Schema)", "10,000 rows × 14 cols\nUCI ML Repository format\nStatistical distributions\nmatching real dataset", BLUE,   0.2),
    ("data_ingestion.py\n(Mapping)",  "Step 1: Generate UCI CSV\nStep 2: Map Type→Equipment\nStep 3: Add medical context\nStep 4: Save both CSVs",  NAVY,   3.55),
    ("medical_equipment\n_maintenance.csv", "10,000 rows × 32 cols\nAll UCI cols retained\n+ 18 medical context cols\nRAG-ready with full text", GREEN,  6.9),
    ("RAG Index\n(Cached)",      "FAISS: 10k × 384-dim\nBM25: tokenised corpus\nEmbeddings: .npy file\nDocuments: .pkl cache",  PURPLE, 10.25),
]
for title, desc, col, x in flow:
    box(sl, x, 1.1, 3.0, 1.9, col)
    txt(sl, title, x + 0.14, 1.18, 2.72, 0.55, size=10, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, desc,  x + 0.14, 1.76, 2.72, 1.18, size=8, color=WHITE)
    if x < 10.25:
        txt(sl, "→", x + 3.05, 1.85, 0.38, 0.36, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# UCI Schema table
box(sl, 0.2, 3.2, 6.2, 4.1, WHITE, BLUE)
txt(sl, "ai4i2020.csv — Exact UCI Schema (14 columns)", 0.35, 3.28, 5.9, 0.34, size=10, bold=True, color=BLUE)
uci = [
    ("UDI",                    "Record index 1–10,000"),
    ("Product ID",             "L/M/H prefix + 5-digit num"),
    ("Type",                   "L=50%  M=30%  H=20%"),
    ("Air temperature [K]",    "~300K, std 2K"),
    ("Process temperature [K]","Air temp + ~10K"),
    ("Rotational speed [rpm]", "~1500 rpm, std 200"),
    ("Torque [Nm]",            "~40 Nm, std 10"),
    ("Tool wear [min]",        "0–253 min (uniform)"),
    ("Machine failure",        "1 = failed (~1.9%)"),
    ("TWF",                    "Tool Wear Failure"),
    ("HDF",                    "Heat Dissipation Failure"),
    ("PWF",                    "Power Failure"),
    ("OSF",                    "Overstrain Failure"),
    ("RNF",                    "Random Failure (~0.1%)"),
]
cy = 3.7
cx = 0.38
for i, (col, desc) in enumerate(uci):
    col_x = cx + (i % 2) * 3.0
    col_y = cy + (i // 2) * 0.42
    txt(sl, col + ":", col_x, col_y, 1.55, 0.32, size=7.5, bold=True, color=BLUE)
    txt(sl, desc,       col_x + 1.57, col_y, 1.38, 0.32, size=7.5, color=DGRAY)

# Mapping table
box(sl, 6.6, 3.2, 6.5, 4.1, WHITE, GREEN)
txt(sl, "Medical Equipment Mapping  (data_ingestion.py)", 6.75, 3.28, 6.2, 0.34, size=10, bold=True, color=GREEN)
mapping = [
    ("Type L (50%)",              "→ Ventilator, Infusion Pump, Patient Monitor"),
    ("Type M (30%)",              "→ Ultrasound, ICU Device, Lab Analyzer"),
    ("Type H (20%)",              "→ MRI, CT Scanner"),
    ("TWF / HDF / PWF / OSF / RNF","→ failure_type (human-readable string)"),
    ("UDI mod len(options)",      "→ equipment_id  (e.g. MRI001, VEN002)"),
    ("Equipment type",            "→ hospital_unit  (e.g. Radiology, ICU)"),
    ("Machine failure + severity","→ downtime_hours, maintenance_cost_usd"),
    ("Technician notes",          "5 categories × 4 templates each"),
    ("Timestamps",                "2023-01-01 → 2024-12-31  (random)"),
    ("Total output cols",         "32 columns  (14 UCI + 18 enriched)"),
]
my = 3.7
for src, tgt in mapping:
    txt(sl, src + ":", 6.75, my, 2.45, 0.3, size=7.5, bold=True, color=GREEN)
    txt(sl, tgt,       9.22, my, 3.72, 0.3, size=7.5, color=DGRAY)
    my += 0.4


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Technology Stack
# ══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
bg(sl, LGRAY)
box(sl, 0, 0, 13.33, 0.65, NAVY)
txt(sl, "Technology Stack", 0.35, 0.08, 10, 0.48, size=22, bold=True)
txt(sl, "All libraries, versions and roles in the system",
    0.35, 0.38, 10, 0.28, size=10, color=CYAN)

categories = [
    ("FRONTEND", BLUE, [
        ("React 18",          "UI component framework"),
        ("Vite 8.0",          "Build tool & dev server"),
        ("Recharts",          "BarChart, LineChart, PieChart"),
        ("Axios",             "HTTP client for API calls"),
        ("CSS Grid / Flex",   "Responsive layout system"),
    ]),
    ("BACKEND", NAVY, [
        ("FastAPI 0.104",     "REST API framework"),
        ("Uvicorn",           "ASGI server"),
        ("Pydantic v2",       "Request / response models"),
        ("pandas 2.1",        "DataFrame operations"),
        ("numpy 1.26",        "Embedding arithmetic"),
    ]),
    ("AI / LLM", PURPLE, [
        ("Claude Sonnet 4.6", "4 AI analysis agents"),
        ("GPT-4o mini",       "Input + output guardrails"),
        ("anthropic ≥0.25",   "Anthropic API client"),
        ("openai ≥2.0",       "OpenAI API client"),
        ("LangGraph 1.2",     "StateGraph multi-agent pipeline"),
    ]),
    ("RAG / SEARCH", GREEN, [
        ("FAISS-CPU 1.7.4",          "Vector similarity index"),
        ("sentence-transformers",     "all-MiniLM-L6-v2 embeddings"),
        ("rank-bm25 0.2.2",          "BM25Okapi keyword search"),
        ("scikit-learn 1.3",         "ML utilities"),
        ("numpy (embeddings.npy)",   "Reranking embedding cache"),
    ]),
    ("EVALUATION", AMBER, [
        ("DeepEval 4.0.5",     "RAG quality metrics library"),
        ("LLM-as-Judge",       "4-dimension rubric scoring"),
        ("AnswerRelevancy",    "DeepEval metric"),
        ("Faithfulness",       "DeepEval metric"),
        ("Contextual metrics", "Precision + Recall"),
    ]),
    ("DATASET", CYAN, [
        ("ai4i2020.csv",              "UCI AI4I 2020 (10k rows)"),
        ("medical_equip_maint.csv",   "Enriched (32 cols)"),
        ("faiss_index.bin",           "15 MB FAISS index"),
        ("embeddings.npy",            "10k × 384-dim float32"),
        ("bm25.pkl + documents.pkl",  "BM25 + doc text cache"),
    ]),
]

for i, (cat, col, items) in enumerate(categories):
    row   = i // 3
    col_i = i % 3
    bx = 0.2 + col_i * 4.35
    by = 0.82 + row * 3.28
    box(sl, bx, by, 4.1, 3.1, WHITE, col)
    box(sl, bx, by, 4.1, 0.44, col)
    txt(sl, cat, bx + 0.15, by + 0.07, 3.8, 0.34, size=10, bold=True)
    iy = by + 0.54
    for lib, purpose in items:
        box(sl, bx + 0.15, iy, 3.8, 0.42, LGRAY)
        txt(sl, lib,     bx + 0.28, iy + 0.05, 1.65, 0.32, size=8,   bold=True, color=col)
        txt(sl, purpose, bx + 1.98, iy + 0.05, 1.95, 0.32, size=7.5, color=DGRAY)
        iy += 0.48


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — End-to-End Data Flow (step by step)
# ══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
bg(sl, DGRAY)
box(sl, 0, 0, 13.33, 0.65, CYAN)
txt(sl, "End-to-End Data Flow — User Query to Recommendation", 0.35, 0.08, 11, 0.48, size=21, bold=True, color=NAVY)

steps = [
    ("1", "User types query in React AI Assistant chat",
     "e.g.  'Critical MRI overheating failure in Radiology unit — recommend maintenance steps'",    BLUE),
    ("2", "Axios POST /api/query  →  FastAPI endpoint",
     "Payload: { query, k=5, alpha=0.6, equipment_type?, hospital_unit?, severity?, failure_only? }", CYAN),
    ("3", "FastAPI calls run_langgraph_pipeline()  →  pipeline.invoke(MedEquipState)",
     "Initial state: 17 fields initialised — query, filters, k, alpha; is_valid=False; logs=[]",     NAVY),
    ("4", "node_input_guardrail()  →  GPT-4o mini / rule-based",
     "Returns: is_valid=True/False, intent='failure_analysis', urgency='immediate', refined_query='…'", AMBER),
    ("5", "node_rag_retrieval()  →  Hybrid FAISS + BM25 + Embedding Rerank",
     "Returns: retrieved_docs[k]  each with  relevance_score, document_text, equipment_id, failure_type, …", GREEN),
    ("6", "4 × Agent nodes  →  Claude Sonnet 4.6 / rule-based  (sequential via LangGraph edges)",
     "Retrieval → Reliability → Maintenance → Recommendation.  Each reads state and writes its result back.", PURPLE),
    ("7", "node_output_guardrail()  →  GPT-4o mini / rule-based",
     "Returns: quality_score (0-100), safety_rating='safe/caution/unsafe', approved_recommendation",    AMBER),
    ("8", "FastAPI maps final LangGraph state → QueryResponse JSON",
     "agent_analysis, guardrails{input+output}, pipeline_log[], node_timings{}, refined_query, processing_time_ms", BLUE),
    ("9", "React renders AssistantMessage component",
     "GuardrailBadges + RefinedQueryNote + PipelineLog + KeyFindings + ActionPlan + ScoreCircle",        CYAN),
]

for i, (num, title, detail, col) in enumerate(steps):
    y = 0.82 + i * 0.72
    box(sl, 0.2, y, 0.48, 0.56, col)
    txt(sl, num,    0.2,  y + 0.08, 0.48, 0.4, size=14, bold=True, align=PP_ALIGN.CENTER)
    box(sl, 0.78,  y, 12.3, 0.56, RGBColor(0x26, 0x37, 0x50), col)
    txt(sl, title,  0.94, y + 0.04, 5.2,  0.26, size=9,   bold=True, color=col)
    txt(sl, detail, 0.94, y + 0.3,  12.0, 0.24, size=7.5, color=RGBColor(0xb8, 0xca, 0xdf))


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Requirement Coverage Summary
# ══════════════════════════════════════════════════════════════════════════
sl = blank_slide()
bg(sl, LGRAY)
box(sl, 0, 0, 13.33, 0.65, NAVY)
txt(sl, "Requirement Coverage — All Items Implemented", 0.35, 0.08, 10, 0.48, size=22, bold=True)
txt(sl, "Requirement 1 (Basic)  +  Requirement 2 (Advanced)  +  Multi-Agent System  +  Frontend — fully implemented & verified",
    0.35, 0.38, 12, 0.28, size=9.5, color=CYAN)

# Req 1
box(sl, 0.2, 0.82, 6.1, 6.0, WHITE, BLUE)
txt(sl, "✅  REQUIREMENT 1 — Basic", 0.35, 0.9, 5.8, 0.36, size=11, bold=True, color=BLUE)
req1 = [
    ("Basic RAG incident retrieval",   "FAISS IndexFlatIP · all-MiniLM-L6-v2 · 10,000 records"),
    ("Hybrid search",                  "FAISS α=0.6 + BM25Okapi α=0.4 · normalised score fusion"),
    ("Semantic search",                "384-dim sentence embeddings · cosine similarity · cached"),
    ("Metadata filtering",             "equipment_type, hospital_unit, severity, failure_only"),
    ("Root-cause recommendation",      "recommendation_agent synthesises all 4 agent outputs"),
    ("Input validation guardrails",    "GPT-4o mini + rule-based fallback in guardrails.py"),
    ("Incident similarity ranking",    "3-phase: hybrid → embedding rerank → metadata filter"),
    ("Maintenance recommendations",    "maintenance_agent · prioritised action plans + urgency"),
    ("API endpoints for integrations", "12 REST endpoints via FastAPI with Pydantic validation"),
]
ry = 1.36
for req, impl in req1:
    box(sl, 0.35, ry, 5.8, 0.54, LGRAY, BLUE)
    txt(sl, "✓  " + req, 0.5, ry + 0.04, 5.55, 0.24, size=8.5, bold=True, color=BLUE)
    txt(sl, impl,         0.5, ry + 0.28, 5.55, 0.24, size=7.5, color=GRAY)
    ry += 0.59

# Req 2
box(sl, 6.6, 0.82, 6.5, 6.0, WHITE, PURPLE)
txt(sl, "✅  REQUIREMENT 2 — Advanced", 6.75, 0.9, 6.2, 0.36, size=11, bold=True, color=PURPLE)
req2 = [
    ("DeepEval evaluation",          "evaluation.py · 4 metrics: AnswerRelevancy, Faithfulness,\nContextualPrecision, ContextualRecall · /api/evaluate"),
    ("Anomaly correlation analysis",  "/api/analyze/anomaly · temp + wear + RPM\nthreshold anomaly detection with correlation report"),
    ("Reranking (maintenance embs)", "rag_system.rerank() · cosine_sim + metadata bonus\n+ recency · 3-phase pipeline in hybrid_search()"),
    ("LLM-as-judge validation",      "llm_judge.py · 4-dim rubric 0–100\n(Safety/Accuracy/Actionability/Evidence) · /api/judge"),
    ("Token optimization",           "_approx_tokens() word-count budget · max_tokens per\nagent (800-1500) · context truncation in get_context()"),
    ("Equipment Retrieval Agent",    "Identifies relevant equipment + query intent"),
    ("Reliability Analysis Agent",   "Failure rate, reliability score 0-100, anomalies"),
    ("Maintenance Planning Agent",   "Prioritised actions, urgency, preventive measures"),
    ("Recommendation Agent",         "Final synthesis, confidence score, root cause"),
    ("Frontend interface",           "6-view React app · chat, dashboard, evaluation,\npipeline, anomaly detection, maintenance records"),
]
ry = 1.36
for req, impl in req2:
    box(sl, 6.75, ry, 6.2, 0.54, LGRAY, PURPLE)
    txt(sl, "✓  " + req, 6.9, ry + 0.04, 5.95, 0.24, size=8.5, bold=True, color=PURPLE)
    txt(sl, impl,         6.9, ry + 0.28, 5.95, 0.24, size=7.5, color=GRAY)
    ry += 0.54

# Footer
box(sl, 0.2, 6.95, 12.9, 0.42, NAVY)
txt(sl, "Stack: FastAPI · LangGraph 1.2 · DeepEval 4.0 · FAISS · BM25 · sentence-transformers · Claude Sonnet 4.6 · GPT-4o mini · React 18 · UCI AI4I 2020 (10,000 records)",
    0.4, 7.0, 12.5, 0.32, size=8, color=WHITE, align=PP_ALIGN.CENTER)


# ─── Save ──────────────────────────────────────────────────────────────────
out = "/home/labuser/capstone med/MedEquip_AI_Architecture.pptx"
prs.save(out)
print(f"Saved → {out}")
